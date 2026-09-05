#!/usr/bin/env python3
"""Bounded, non-executing document extraction sidecar for WellForgeRag."""

from __future__ import annotations

import json
import sys
from pathlib import Path
from typing import Any

MAX_TEXT_CHARS = 8_000_000
MAX_TABLE_ROWS = 500
MAX_TABLE_COLUMNS = 128


def section(locator: str, text: str) -> dict[str, str]:
    return {"locator": locator, "text": text[:MAX_TEXT_CHARS]}


def column(name: str, data_type: str | None = None) -> dict[str, Any]:
    return {"name": name, "data_type": data_type, "nullable": None}


def envelope(backend: str) -> dict[str, Any]:
    return {
        "family": "document",
        "status": "extracted",
        "backend": backend,
        "text_sections": [],
        "profiles": [],
        "metadata": {},
        "warnings": [],
    }


def extract_pdf(path: Path) -> dict[str, Any]:
    import pdfplumber

    result = envelope("python-pdfplumber")
    with pdfplumber.open(path) as pdf:
        result["metadata"]["pages"] = len(pdf.pages)
        for index, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                result["text_sections"].append(section(f"page:{index}", text))
            for table_index, table in enumerate(page.extract_tables() or [], start=1):
                rows = table[:MAX_TABLE_ROWS]
                width = max((len(row or []) for row in rows), default=0)
                result["profiles"].append(
                    {
                        "name": f"page-{index}-table-{table_index}",
                        "row_count": len(table),
                        "columns": [column(f"column_{i + 1}") for i in range(min(width, MAX_TABLE_COLUMNS))],
                    }
                )
                sample = "\n".join(
                    "\t".join("" if cell is None else str(cell) for cell in (row or [])[:MAX_TABLE_COLUMNS])
                    for row in rows
                )
                if sample.strip():
                    result["text_sections"].append(
                        section(f"page:{index}:table:{table_index}:sample", sample)
                    )
    return result


def extract_docx(path: Path) -> dict[str, Any]:
    import docx

    document = docx.Document(path)
    result = envelope("python-docx")
    paragraphs = [p.text for p in document.paragraphs if p.text.strip()]
    if paragraphs:
        result["text_sections"].append(section("document:paragraphs", "\n\n".join(paragraphs)))
    for index, table in enumerate(document.tables, start=1):
        rows = list(table.rows)
        width = max((len(row.cells) for row in rows), default=0)
        result["profiles"].append(
            {
                "name": f"table-{index}",
                "row_count": len(rows),
                "columns": [column(f"column_{i + 1}") for i in range(min(width, MAX_TABLE_COLUMNS))],
            }
        )
        sample_rows = rows[:MAX_TABLE_ROWS]
        text = "\n".join(
            "\t".join(cell.text for cell in row.cells[:MAX_TABLE_COLUMNS]) for row in sample_rows
        )
        if text.strip():
            result["text_sections"].append(section(f"table:{index}:sample", text))
    return result


def extract_pptx(path: Path) -> dict[str, Any]:
    import pptx

    presentation = pptx.Presentation(path)
    result = envelope("python-pptx")
    result["metadata"]["slides"] = len(presentation.slides)
    for index, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                parts.append(text)
        if parts:
            result["text_sections"].append(section(f"slide:{index}", "\n".join(parts)))
    return result


def _sheet_result(result: dict[str, Any], name: str, rows: int | None, width: int, samples: list[list[str]]) -> None:
    result["profiles"].append(
        {
            "name": name,
            "row_count": rows,
            "columns": [column(f"column_{i + 1}") for i in range(min(width, MAX_TABLE_COLUMNS))],
        }
    )
    rendered = "\n".join("\t".join(row[:MAX_TABLE_COLUMNS]) for row in samples)
    if rendered.strip():
        result["text_sections"].append(section(f"sheet:{name}:sample", rendered))


def extract_xlsx(path: Path) -> dict[str, Any]:
    import openpyxl

    workbook = openpyxl.load_workbook(path, data_only=False, read_only=True, keep_links=False)
    result = envelope("python-openpyxl")
    result["family"] = "structured-data"
    result["metadata"]["sheets"] = list(workbook.sheetnames)
    for name in workbook.sheetnames:
        worksheet = workbook[name]
        samples: list[list[str]] = []
        width = 0
        for row_index, row in enumerate(worksheet.iter_rows(values_only=True), start=1):
            values = ["" if value is None else str(value) for value in row]
            width = max(width, len(values))
            if row_index <= MAX_TABLE_ROWS and any(value.strip() for value in values):
                samples.append(values)
            if row_index >= MAX_TABLE_ROWS and worksheet.max_row is None:
                break
        _sheet_result(result, name, worksheet.max_row, width, samples)
    workbook.close()
    return result


def extract_xls(path: Path) -> dict[str, Any]:
    import xlrd

    workbook = xlrd.open_workbook(path, on_demand=True)
    result = envelope("python-xlrd")
    result["family"] = "structured-data"
    result["metadata"]["sheets"] = workbook.sheet_names()
    for sheet in workbook.sheets():
        samples: list[list[str]] = []
        for row_index in range(min(sheet.nrows, MAX_TABLE_ROWS)):
            samples.append([str(value) for value in sheet.row_values(row_index)[:MAX_TABLE_COLUMNS]])
        _sheet_result(result, sheet.name, sheet.nrows, sheet.ncols, samples)
    workbook.release_resources()
    return result


def extract_xlsb(path: Path) -> dict[str, Any]:
    import pyxlsb

    result = envelope("python-pyxlsb")
    result["family"] = "structured-data"
    with pyxlsb.open_workbook(str(path)) as workbook:
        result["metadata"]["sheets"] = list(workbook.sheets)
        for name in workbook.sheets:
            samples: list[list[str]] = []
            row_count = 0
            width = 0
            with workbook.get_sheet(name) as sheet:
                for row in sheet.rows():
                    row_count += 1
                    values = ["" if cell.v is None else str(cell.v) for cell in row]
                    width = max(width, len(values))
                    if row_count <= MAX_TABLE_ROWS and any(value.strip() for value in values):
                        samples.append(values)
            _sheet_result(result, name, row_count, width, samples)
    return result


def extract_rtf(path: Path) -> dict[str, Any]:
    from striprtf.striprtf import rtf_to_text

    text = rtf_to_text(path.read_text(encoding="utf-8", errors="replace"))
    result = envelope("python-striprtf")
    result["text_sections"].append(section("document", text))
    return result


def extract_image(path: Path) -> dict[str, Any]:
    from PIL import Image
    import pytesseract

    result = envelope("python-tesseract")
    with Image.open(path) as image:
        result["metadata"].update({"width": image.width, "height": image.height, "mode": image.mode})
        text = pytesseract.image_to_string(image)
    if text.strip():
        result["text_sections"].append(section("image:ocr", text))
    return result


HANDLERS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
    ".xlsx": extract_xlsx,
    ".xlsm": extract_xlsx,
    ".xls": extract_xls,
    ".xlsb": extract_xlsb,
    ".rtf": extract_rtf,
    ".png": extract_image,
    ".jpg": extract_image,
    ".jpeg": extract_image,
    ".tif": extract_image,
    ".tiff": extract_image,
}


def main() -> int:
    if len(sys.argv) != 2:
        print(json.dumps({"error": "usage: extract.py <file-path>"}), file=sys.stderr)
        return 2
    path = Path(sys.argv[1]).resolve()
    if not path.is_file():
        print(json.dumps({"error": f"not a file: {path}"}), file=sys.stderr)
        return 2
    handler = HANDLERS.get(path.suffix.lower())
    if handler is None:
        print(json.dumps({"error": f"unsupported sidecar extension: {path.suffix.lower()}"}), file=sys.stderr)
        return 3
    try:
        result = handler(path)
        encoded = json.dumps(result, ensure_ascii=False, separators=(",", ":"))
        if len(encoded.encode("utf-8")) > 64 * 1024 * 1024:
            raise ValueError("extraction response exceeds 64 MiB safety ceiling")
        sys.stdout.write(encoded)
        sys.stdout.write("\n")
        return 0
    except Exception as error:
        print(json.dumps({"error": f"{type(error).__name__}: {error}"}), file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
