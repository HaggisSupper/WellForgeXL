#!/usr/bin/env python3
"""Extract text and structure from common document formats."""

import json
import os
import sys
import traceback
from pathlib import Path


def extract_pdf(path: str):
    import pdfplumber
    text_parts = []
    tables = []
    pages = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            page_text = page.extract_text() or ""
            if page_text.strip():
                text_parts.append(page_text)
            page_tables = page.extract_tables() or []
            for table in page_tables:
                tables.append({"page": i, "rows": table})
            pages.append({"page": i, "chars": len(page_text)})
    return {
        "text": "\n\n".join(text_parts),
        "pages": pages,
        "tables": tables,
        "images": [],
        "coordinates": [],
        "warnings": [],
    }


def extract_docx(path: str):
    import docx
    doc = docx.Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    tables = []
    for i, table in enumerate(doc.tables, start=1):
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        tables.append({"index": i, "rows": rows})
    return {
        "text": "\n\n".join(paragraphs),
        "pages": [],
        "tables": tables,
        "images": [],
        "coordinates": [],
        "warnings": [],
    }


def extract_xlsx(path: str):
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    parts = []
    tables = []
    for name in wb.sheetnames:
        ws = wb[name]
        sheet_rows = []
        for row in ws.iter_rows(values_only=True):
            sheet_rows.append([str(cell) if cell is not None else "" for cell in row])
        if sheet_rows:
            tables.append({"sheet": name, "rows": sheet_rows})
            non_empty = ["\t".join(row) for row in sheet_rows if any(cell.strip() for cell in row)]
            if non_empty:
                parts.append(f"Sheet: {name}\n" + "\n".join(non_empty))
    return {
        "text": "\n\n".join(parts),
        "pages": [],
        "tables": tables,
        "images": [],
        "coordinates": [],
        "warnings": [],
    }


def extract_xlsb(path: str):
    import pyxlsb
    parts = []
    tables = []
    with pyxlsb.open_workbook(path) as wb:
        for sheet_name in wb.sheets:
            sheet_rows = []
            with wb.get_sheet(sheet_name) as sheet:
                for row in sheet.rows():
                    sheet_rows.append([str(item.v) if item.v is not None else "" for item in row])
            if sheet_rows:
                tables.append({"sheet": sheet_name, "rows": sheet_rows})
                non_empty = ["\t".join(row) for row in sheet_rows if any(cell.strip() for cell in row)]
                if non_empty:
                    parts.append(f"Sheet: {sheet_name}\n" + "\n".join(non_empty))
    return {
        "text": "\n\n".join(parts),
        "pages": [],
        "tables": tables,
        "images": [],
        "coordinates": [],
        "warnings": [],
    }


def extract_pptx(path: str):
    import pptx
    prs = pptx.Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        if slide_text:
            parts.append(f"Slide {i}:\n" + "\n".join(slide_text))
    return {
        "text": "\n\n".join(parts),
        "pages": [{"page": i} for i in range(1, len(prs.slides) + 1)],
        "tables": [],
        "images": [],
        "coordinates": [],
        "warnings": [],
    }


def extract_txt(path: str):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return {
        "text": text,
        "pages": [],
        "tables": [],
        "images": [],
        "coordinates": [],
        "warnings": [],
    }


def extract_image(path: str):
    try:
        import pytesseract
        from PIL import Image
        text = pytesseract.image_to_string(Image.open(path))
        return {
            "text": text,
            "pages": [],
            "tables": [],
            "images": [],
            "coordinates": [],
            "warnings": [],
        }
    except Exception as e:
        return {
            "text": "",
            "pages": [],
            "tables": [],
            "images": [],
            "coordinates": [],
            "warnings": [f"ocr unavailable: {e}"],
        }


HANDLERS = {
    ".pdf": extract_pdf,
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
    ".xlsm": extract_xlsx,
    ".xlsb": extract_xlsb,
    ".pptx": extract_pptx,
    ".txt": extract_txt,
    ".html": extract_txt,
    ".htm": extract_txt,
    ".png": extract_image,
    ".jpg": extract_image,
    ".jpeg": extract_image,
    ".tif": extract_image,
    ".tiff": extract_image,
}


def main():
    if sys.platform == "win32":
        import io
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
        sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")

    if len(sys.argv) < 2:
        print(json.dumps({"error": "usage: extract.py <file-path>"}), file=sys.stderr)
        sys.exit(1)

    file_path = sys.argv[1]
    ext = Path(file_path).suffix.lower()
    handler = HANDLERS.get(ext)

    if not handler:
        result = {
            "text": "",
            "pages": [],
            "tables": [],
            "images": [],
            "coordinates": [],
            "warnings": [f"no handler for extension {ext}"],
        }
    else:
        try:
            result = handler(file_path)
            result["backend"] = "python-" + handler.__name__.replace("extract_", "")
        except Exception as e:
            result = {
                "text": "",
                "pages": [],
                "tables": [],
                "images": [],
                "coordinates": [],
                "warnings": [f"extraction error: {e}", traceback.format_exc()],
            }

    result["filePath"] = file_path
    print(json.dumps(result, ensure_ascii=False))


if __name__ == "__main__":
    main()
